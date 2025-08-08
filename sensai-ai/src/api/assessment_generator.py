# sensai-ai/src/api/assessment_generator.py

import os
import json
from pathlib import Path
from abc import ABC, abstractmethod
from typing import List
import uuid
import csv

# File Parsing & Writing
import fitz
from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.util import Pt

# Core AI and Data Libraries
import openai
from thefuzz import process as fuzzy_process

# Rich UI/UX
from rich.console import Console
from rich.panel import Panel

class DocumentHandler(ABC):
    def __init__(self, file_path: Path): self.file_path = file_path
    @abstractmethod
    def parse(self) -> list[dict]: pass
    @abstractmethod
    def write_annotated(self, output_path: Path, items_to_highlight: list[dict]): pass

class PDFHandler(DocumentHandler):
    SKY_BLUE = (0.529, 0.808, 0.922); YELLOW = (1, 1, 0)
    def parse(self) -> list[dict]:
        chunks = []
        with fitz.open(self.file_path) as doc:
            for page_num, page in enumerate(doc):
                blocks = page.get_text("blocks", sort=True)
                for block in blocks:
                    if block[6] == 0:
                        raw_text = block[4].replace('\n', ' ').strip()
                        text = raw_text.encode('utf-8', 'ignore').decode('utf-8')
                        if text: chunks.append({"text": text, "location": {"page_number": page_num + 1, "bbox": block[:4]}})
        return chunks
    def write_annotated(self, output_path: Path, items_to_highlight: list[dict]):
        doc = fitz.open(self.file_path)
        for item in items_to_highlight:
            loc, q_type, q_id, answer = item.get("location"), item.get("q_type"), item.get("q_id"), item.get("answer")
            if not all([loc, q_type, q_id, answer, "page_number" in loc]): continue
            page = doc.load_page(loc["page_number"] - 1)
            color = self.SKY_BLUE if q_type == "MCQ" else self.YELLOW
            areas = page.search_for(answer, clip=fitz.Rect(loc["bbox"]))
            if not areas: areas = page.search_for(item["text"], clip=fitz.Rect(loc["bbox"]))
            for area in areas:
                highlight = page.add_highlight_annot(area); highlight.set_colors(stroke=color); highlight.update()
                page.insert_text(area.tl + (-15, -5), f"[Q{q_id}]", fontsize=8, color=(0, 0, 0))
        doc.save(output_path, garbage=4, deflate=True)

class WordHandler(DocumentHandler):
    SKY_BLUE = WD_COLOR_INDEX.TURQUOISE; YELLOW = WD_COLOR_INDEX.YELLOW
    def parse(self) -> list[dict]:
        chunks = []; doc = Document(self.file_path)
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                text = para.text.encode('utf-8', 'ignore').decode('utf-8')
                chunks.append({"text": text, "location": {"paragraph_number": i}})
        return chunks
    def write_annotated(self, output_path: Path, items_to_highlight: list[dict]):
        doc = Document(self.file_path)
        highlights = {}
        for item in items_to_highlight:
            loc = item.get("location")
            if loc and "paragraph_number" in loc:
                p_idx = loc["paragraph_number"]
                if p_idx not in highlights: highlights[p_idx] = []
                highlights[p_idx].append(item)
        for i, para in enumerate(doc.paragraphs):
            if i in highlights:
                items = highlights[i]; color = self.SKY_BLUE if items[0]["q_type"] == "MCQ" else self.YELLOW
                ids = ", ".join([f"Q{item['q_id']}" for item in items])
                id_para = para.insert_paragraph_before(f"[{ids}]")
                id_para.runs[0].font.color.rgb = PptxRGBColor(255, 0, 0); id_para.runs[0].font.bold = True
                for run in para.runs: run.font.highlight_color = color
        doc.save(output_path)

class PowerPointHandler(DocumentHandler):
    SKY_BLUE = PptxRGBColor(135, 206, 235); YELLOW = PptxRGBColor(255, 255, 0)
    def parse(self) -> list[dict]:
        chunks = []
        pres = Presentation(self.file_path)
        for slide_num, slide in enumerate(pres.slides, 1):
            for shape_num, shape in enumerate(slide.shapes):
                if shape.has_text_frame and shape.text.strip():
                    for para_num, para in enumerate(shape.text_frame.paragraphs):
                        if para.text.strip():
                            text = para.text.encode('utf-8', 'ignore').decode('utf-8')
                            chunks.append({"text": text, "location": {"slide_number": slide_num, "shape_index": shape_num, "paragraph_index": para_num}})
        return chunks
    def write_annotated(self, output_path: Path, items_to_highlight: list[dict]):
        pres = Presentation(self.file_path)
        highlights_by_para = {}
        for item in items_to_highlight:
            loc = item.get("location")
            if loc and all(k in loc for k in ["slide_number", "shape_index", "paragraph_index"]):
                key = (loc["slide_number"], loc["shape_index"], loc["paragraph_index"])
                if key not in highlights_by_para: highlights_by_para[key] = []
                highlights_by_para[key].append(item)
        for slide_num, slide in enumerate(pres.slides, 1):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    is_shape_annotated = False
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        key = (slide_num, shape_idx, para_idx)
                        if key in highlights_by_para:
                            items = highlights_by_para[key]
                            color = self.SKY_BLUE if items[0]["q_type"] == "MCQ" else self.YELLOW
                            for run in para.runs: run.font.fill.solid(); run.font.fill.fore_color.rgb = color
                            is_shape_annotated = True
                    if is_shape_annotated:
                        shape_items = [item for key, items in highlights_by_para.items() if key[0] == slide_num and key[1] == shape_idx for item in items]
                        q_ids = ", ".join(sorted(list(set([f"Q{item['q_id']}" for item in shape_items]))))
                        left, top, width, height = shape.left, shape.top, shape.width, Pt(20)
                        id_box = slide.shapes.add_textbox(left, top - height if top > height else top, width, height)
                        id_frame = id_box.text_frame; id_frame.text = f"[{q_ids}]"
                        id_para = id_frame.paragraphs[0]; id_para.font.bold = True; id_para.font.size = Pt(10)
                        id_para.font.color.rgb = PptxRGBColor(255, 0, 0)
        pres.save(output_path)

class AssessmentGenerator:
    API_KEY = "sk-acTVFF7PJXIMfYxx-JTVjA"
    MODEL = "openai/gpt-4o-mini"
    CONTEXT_LIMIT = 100000
    
    def __init__(self, output_dir: str = "chatbot"):
        self.console = Console(); self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.client = openai.OpenAI(api_key=self.API_KEY, base_url="https://agent.dev.hyperverge.org")
        self.parsed_chunks = []

    def _get_handler(self, file_path: Path) -> DocumentHandler:
        ext = file_path.suffix.lower()
        if ext == ".pdf": return PDFHandler(file_path)
        if ext == ".docx": return WordHandler(file_path)
        if ext == ".pptx": return PowerPointHandler(file_path)
        raise ValueError(f"Unsupported file type: {ext}")

    def _call_llm(self, prompt: str) -> dict | None:
        try:
            response = self.client.chat.completions.create(
                model=self.MODEL, messages=[{"role": "system", "content": "You are a world-class AI assistant. Follow instructions precisely and adhere to the requested JSON format."}, {"role": "user", "content": prompt}],
                response_format={"type": "json_object"}, temperature=0.2)
            content = response.choices[0].message.content
            if isinstance(content, str): return json.loads(content)
            if isinstance(content, dict): return content
            return None
        except Exception as e:
            self.console.print(f"[bold red]API call failed. The hardcoded API key is likely invalid. Error: {e}[/bold red]")
            return None

    def _generate_question_type(self, q_type: str, count: int, full_text: str) -> List[dict]:
        self.console.print(f"[yellow]Stage: Generating {count} {q_type}s...[/yellow]")
        q_name = "Multiple-Choice Questions" if q_type == "MCQ" else "Short-Answer Questions"
        json_key = "multiple_choice_questions" if q_type == "MCQ" else "short_answer_questions"
        opts = '"options": (list of 4 strings)' if q_type == "MCQ" else '"options" key should be omitted'
        prompt = f'Based on the document text, generate an assessment with exactly {count} {q_name}. Respond with a single JSON object with ONE key: "{json_key}". Each question must contain: "question", "answer", "hint", "source_quote", and {opts}. DOCUMENT TEXT: ---\n{full_text}\n---'
        generated_qs = []
        if len(full_text) < self.CONTEXT_LIMIT:
            result = self._call_llm(prompt)
            if result: generated_qs = result.get(json_key, [])
        return generated_qs[:count]

    def run(self, input_path_str: str):
        input_path = Path(input_path_str)
        self.console.print(Panel(f"🚀 Starting Simplified Pipeline for [cyan]{input_path.name}[/cyan] 🚀", border_style="bold magenta"))
        handler = self._get_handler(input_path)
        self.parsed_chunks = handler.parse()
        
        full_text = "\n\n".join([c['text'] for c in self.parsed_chunks])
        
        mcq_list = self._generate_question_type("MCQ", 15, full_text)
        saq_list = self._generate_question_type("SAQ", 5, full_text)

        self.console.print("[yellow]Finalizing questions and preparing annotated document...[/yellow]")
        final_mcqs, final_saqs, items_to_highlight = [], [], []
        q_id_counter = 1
        
        combined_list = (mcq_list or []) + (saq_list or [])

        for q_object in combined_list:
            try:
                quote, answer = q_object.get("source_quote"), q_object.get("answer")
                if not quote or not answer: continue
                
                match_results = fuzzy_process.extract(quote, [c['text'] for c in self.parsed_chunks], limit=1)
                if not match_results or match_results[0][1] < 75: continue
                
                best_match_text, score = match_results[0]
                chunk = next((c for c in self.parsed_chunks if c['text'] == best_match_text), None)
                
                if chunk:
                    q_object["id"] = q_id_counter
                    q_object["citation"] = chunk["location"]
                    q_type = "MCQ" if "options" in q_object else "SAQ"
                    items_to_highlight.append({"q_id": q_id_counter, "q_type": q_type, "text": chunk["text"], "answer": str(answer), "location": chunk["location"]})
                    
                    if q_type == "MCQ":
                        final_mcqs.append(q_object)
                    else:
                        final_saqs.append(q_object)
                    
                    q_id_counter += 1
            except Exception as e: 
                self.console.print(f"[bold red]Skipping one question due to error: {e}[/bold red]")

        final_assessment = {"multiple_choice_questions": final_mcqs, "short_answer_questions": final_saqs}

        annotated_filename = f"highlighted_{input_path.stem}{input_path.suffix}"
        annotated_path = self.output_dir / annotated_filename
        handler.write_annotated(annotated_path, items_to_highlight)

        unique_id = str(uuid.uuid4())
        json_filename = f"assessment_{unique_id}.json"
        json_path = self.output_dir / json_filename
        with open(json_path, 'w', encoding='utf-8') as f: json.dump(final_assessment, f, indent=2)

        self.console.print(Panel(f"✅ [bold green]SUCCESS! Your files are in the '{self.output_dir.name}' folder.[/bold green] ✅", title="Execution Complete"))

        return {
            "json": str(json_path),
            "annotated_doc": str(annotated_path),
        }

# --- Scoring logic has been moved to assessment.py for more direct control ---